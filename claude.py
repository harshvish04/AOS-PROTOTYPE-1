 from cmath import e
from email.mime import audio

from tkinter.messagebox import YES
import pyttsx3
import speech_recognition as sr 
import datetime
import wikipedia
import webbrowser
import pywhatkit
import pptx


engine =  pyttsx3.init('sapi5')
voices = engine.getProperty('voices')

engine.setProperty('voice',voices[1].id)

def wishMe():
   hour=int(datetime.datetime.now().hour)
   if hour>=0 and hour<12:
       speak("Good Morning!")
   elif hour>-12 and hour<18:
       speak("Good Afternoon")
   else:
       speak("Good Evening!")
       
   speak("I am cloudiee Sir Please tell me how mayIhelp you")
def takeCommand():
    #It takes microphone input from the user and returns string output
   r=sr.Recognizer()
   
   with sr.Microphone() as source:
       print("Listening ...")  
    
       audio= r.listen(source)
       
       print("Rcognizing")  
       query = r.recognize_google(audio, language='en-in') 
       print(f"User said: {query}\n")  
        
       
       return query



   

def speak(audio):
    engine.say(audio)
    engine.runAndWait()

if __name__ =="__main__":
    wishMe()
          
           
def create_presentation(topic):
    prs = pptx.Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = topic

    prs.save(f'{topic}.pptx')           



          
def run_claude():
    command = takeCommand()
    print (command)
    if 'play' in command:
     song = command.replace('play', '')
     speak('playing' + song) 
     pywhatkit.playonyt (song)
    if 'search' in command:
     search = command.replace('search', '')
     speak('searching' + search) 
     pywhatkit.playonyt (search)   
    elif 'time' in command:
      time = datetime.datetime.now().strftime('%H:%M')
      print(time)
      speak('Current time is ' + time)    
    elif 'who is' in command:
        
     ans = command.replace('who is', '')
     
     info = wikipedia.summary(ans, 5)
     print (info)
     speak(info)   
    elif 'what is ' in command:
     speak('I am searching please wait sir')   
     ans = command.replace('what is', '')
    elif 'find the code of ' in command:
     speak('I am searching please wait sir')   
     ans = command.replace('find the code of', '') 
     info = wikipedia.summary(ans, 2)
     speak("Sir i have find that")
     print (info)
     speak(info)
     
   
                 
run_claude()         

